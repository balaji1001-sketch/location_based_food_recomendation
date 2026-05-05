from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
PRIMARY_COLOR = RGBColor(25, 118, 210)
SECONDARY_COLOR = RGBColor(66, 133, 244)
ACCENT_COLOR = RGBColor(251, 188, 5)
TEXT_COLOR = RGBColor(33, 33, 33)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(28)
    subtitle_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(title, bullets):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Add colored line below title
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)
    
    # Add bullet points
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)

def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Add colored line
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = ACCENT_COLOR
    line.line.width = Pt(3)
    
    # Left column
    left_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4), Inches(0.5))
    left_title_frame = left_title_box.text_frame
    left_title_frame.text = left_title
    left_title_frame.paragraphs[0].font.size = Pt(24)
    left_title_frame.paragraphs[0].font.bold = True
    left_title_frame.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    
    left_content = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(4.2), Inches(4.8))
    left_text_frame = left_content.text_frame
    left_text_frame.word_wrap = True
    
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = left_text_frame.paragraphs[0]
        else:
            p = left_text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Right column
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4), Inches(0.5))
    right_title_frame = right_title_box.text_frame
    right_title_frame.text = right_title
    right_title_frame.paragraphs[0].font.size = Pt(24)
    right_title_frame.paragraphs[0].font.bold = True
    right_title_frame.paragraphs[0].font.color.rgb = SECONDARY_COLOR
    
    right_content = slide.shapes.add_textbox(Inches(5.2), Inches(2.2), Inches(4.2), Inches(4.8))
    right_text_frame = right_content.text_frame
    right_text_frame.word_wrap = True
    
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = right_text_frame.paragraphs[0]
        else:
            p = right_text_frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(4)
        p.space_after = Pt(4)

# Slide 1: Title Slide
add_title_slide("Food Recommendation &", "Price Comparison System")

# Slide 2: Project Overview
add_content_slide("Project Overview", [
    "🎯 Location-based restaurant discovery platform",
    "💰 Real-time price comparison across restaurants",
    "🤖 AI-powered personalized recommendations",
    "⭐ User reviews with sentiment analysis",
    "🗺️ Geospatial search with distance filtering",
    "🔐 Secure authentication with JWT tokens"
])

# Slide 3: Core Features
add_content_slide("Core Features", [
    "✓ Location-Based Search - Find restaurants within specified radius",
    "✓ Advanced Filtering - Cuisine, price range, rating, dietary needs",
    "✓ Personalized Recommendations - Content & collaborative filtering",
    "✓ Price Comparison - Compare dish prices across restaurants",
    "✓ Review System - Submit & view reviews with ratings",
    "✓ Sentiment Analysis - NLP-based review analysis"
])

# Slide 4: Architecture Overview
add_two_column_slide("Architecture", 
    "Frontend",
    ["• React.js 18.2.0", "• Redux state management", "• React Router", "• Axios HTTP client", "• Chart.js visualization", "•Leaflet maps"],
    "Backend",
    ["• Node.js runtime", "• Express.js 4.18.2", "• MongoDB 7.0.0", "• Mongoose ODM", "• Redis caching", "• JWT authentication"]
)

# Slide 5: Tech Stack
add_two_column_slide("Technology Stack",
    "ML & Scraping",
    ["• Python 3.8+", "• Flask 2.3.0", "• scikit-learn", "• NLTK & TextBlob", "• BeautifulSoup4", "• Scrapy 2.9.0"],
    "Security & Data",
    ["• bcryptjs password hashing", "• JWT token auth", "• MongoDB geospatial indexes", "• Redis for caching", "• APScheduler", "• Sentiment analysis"]
)

# Slide 6: Database Models
add_content_slide("Database Models", [
    "📊 Users - User accounts with preferences & ratings",
    "🏪 Restaurants - Restaurant info with geospatial coordinates",
    "💬 Reviews - User reviews with sentiment scores",
    "⭐ UserRatings - Collaborative filtering data",
    "💵 PriceTracking - Historical price tracking",
    "🎯 Key Indexes: Geospatial (2dsphere), Compound indexes, TTL indexes"
])

# Slide 7: APIs - Authentication & Restaurants
add_two_column_slide("API Endpoints (1/3)",
    "Authentication",
    ["POST /auth/register", "POST /auth/login", "GET /auth/me", "PUT /auth/profile"],
    "Restaurants",
    ["GET /restaurants (filters)", "GET /restaurants/:id", "GET /restaurants/search", "GET /restaurants/trending"]
)

# Slide 8: APIs - Reviews & Recommendations
add_two_column_slide("API Endpoints (2/3)",
    "Reviews",
    ["POST /reviews/:id", "GET /reviews/:id", "POST /reviews/:id/helpful"],
    "Recommendations",
    ["GET /recommendations", "POST /recommendations/rate"]
)

# Slide 9: APIs - Price & Search
add_content_slide("API Endpoints (3/3)", [
    "💰 Price Comparison & Trends",
    "  • GET /price-comparison?foodItem=X&days=N",
    "  • GET /price-trends?foodItem=X&days=30",
    "  • POST /price-tracking",
    "",
    "📝 Base URL: http://localhost:5000/api",
    "🔐 Protected endpoints require JWT token"
])

# Slide 10: Recommendation Algorithms
add_content_slide("ML Algorithms", [
    "🎯 Content-Based Filtering",
    "  • Feature vectors: price, rating, cuisine, budget",
    "  • Cosine similarity calculation",
    "",
    "👥 Collaborative Filtering",
    "  • User-user similarity matrix",
    "  • Discovers unexpected recommendations",
    "",
    "🔗 Hybrid Approach - 60% content + 40% collaborative"
])

# Slide 11: Sentiment Analysis
add_content_slide("Sentiment Analysis", [
    "📊 VADER - Rule-based sentiment analyzer",
    "📈 TextBlob - Statistical polarity scores",
    "🤖 Ensemble Method - Combined scores for robustness",
    "",
    "Classifications:",
    "  • Positive sentiment (score > 0.5)",
    "  • Neutral sentiment (-0.5 to 0.5)",
    "  • Negative sentiment (score < -0.5)"
])

# Slide 12: Data Flow
add_content_slide("Data Flow Examples", [
    "🔍 Search Flow:",
    "  User location → Filter criteria → Database query → Results on map",
    "",
    "🤖 Recommendation Flow:",
    "  User behavior → ML analysis → Similar users/items → Recommendations",
    "",
    "💬 Review Flow:",
    "  User review → Save → Sentiment analysis → Update ratings"
])

# Slide 13: Performance Optimization
add_two_column_slide("Optimizations",
    "Backend",
    ["• Geospatial indexes", "• Compound indexes", "• Query optimization", "• Connection pooling"],
    "Frontend",
    ["• Code splitting", "• Lazy loading", "• Image optimization", "• CSS minification"]
)

# Slide 14: Deployment & Setup
add_content_slide("Quick Start", [
    "Prerequisites: Node.js v14+, Python 3.8+, MongoDB 4.4+, Redis 6.0+",
    "",
    "Services:",
    "  • Frontend: http://localhost:3000",
    "  • Backend: http://localhost:5000",
    "  • ML Service: http://localhost:5001",
    "",
    "Installation: npm install → python -m venv → Configure .env → npm/python start"
])

# Slide 15: Key Features Summary
add_content_slide("Key Strengths", [
    "✨ Microservices Architecture - Scalable & maintainable",
    "🧠 Advanced ML Models - Accurate recommendations",
    "⚡ Real-time Performance - Redis caching & indexing",
    "🌍 Geospatial Engine - Precise location-based search",
    "🔒 Security First - JWT auth & password hashing",
    "📱 Responsive UI - Mobile-friendly design"
])

# Slide 16: Future Enhancements
add_content_slide("Future Roadmap", [
    "🚀 Mobile App - Native iOS/Android applications",
    "🤖 Advanced ML - Deep learning recommendations",
    "🔔 Push Notifications - Real-time alerts",
    "💳 Payment Integration - Online ordering",
    "📊 Analytics Dashboard - Restaurant insights",
    "🌐 Multi-language Support - Global expansion"
])

# Slide 17: Closing Slide
add_title_slide("Thank You!", "Questions?")

# Save presentation
prs.save('d:\\project\\Food_Recommendation_System_Presentation.pptx')
print("✓ PowerPoint presentation created successfully!")
print("📁 Saved as: Food_Recommendation_System_Presentation.pptx")
